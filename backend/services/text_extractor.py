"""
Service d'extraction de texte pour l'AI CFO Suite.
Supporte : PDF (natif + OCR), DOCX, XLSX/XLS, CSV, images, texte brut.
Extrait de main.py pour séparation des responsabilités.
"""
import io
import sys
from pathlib import Path

# ─────────────────────────────────────────────────────────────────────────────
# Imports optionnels — chaque librairie peut manquer sans bloquer le reste
# ─────────────────────────────────────────────────────────────────────────────
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    import openpyxl  # noqa: F401 (requis par pandas pour xlsx)
    import pandas as pd
except ImportError:
    pd = None

try:
    import pytesseract
    from pdf2image import convert_from_bytes
    from PIL import Image

    # Chemin Tesseract sous Windows
    if sys.platform == "win32":
        _tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if Path(_tesseract_path).exists():
            pytesseract.pytesseract.tesseract_cmd = _tesseract_path

    # Chemin Poppler sous Windows (dossier local)
    POPPLER_PATH: str | None = None
    if sys.platform == "win32":
        _local_poppler = (
            Path(__file__).parent.parent
            / "poppler"
            / "poppler-24.08.0"
            / "Library"
            / "bin"
        )
        if _local_poppler.exists():
            POPPLER_PATH = str(_local_poppler)

except ImportError:
    Image = None
    pytesseract = None
    convert_from_bytes = None
    POPPLER_PATH = None


# ─────────────────────────────────────────────────────────────────────────────
# Extraction principale
# ─────────────────────────────────────────────────────────────────────────────

def extract_text_from_file(file_content: bytes, filename: str) -> str:
    """
    Extrait le texte brut d'un fichier selon son extension.

    Formats supportés:
    - PDF : extraction native (PyPDF2) + OCR de secours (Tesseract)
    - DOCX : paragraphes (python-docx)
    - XLSX/XLS : toutes les feuilles (pandas + openpyxl)
    - CSV : tableau formaté (pandas)
    - Images (PNG, JPG, BMP, TIFF, GIF) : OCR Tesseract
    - Texte brut (TXT, JSON, XML, MD, LOG) : décodage UTF-8

    Returns:
        Texte extrait, ou message d'erreur préfixé par [ERREUR].
    """
    if not file_content:
        return "[ERREUR] Fichier vide"

    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    try:
        if ext == "pdf":
            return _extract_pdf(file_content, filename)
        elif ext in ("docx", "doc"):
            return _extract_docx(file_content)
        elif ext in ("pptx", "ppt"):
            return _extract_pptx(file_content)
        elif ext in ("xlsx", "xls"):
            return _extract_excel(file_content)
        elif ext == "csv":
            return _extract_csv(file_content)
        elif ext in ("png", "jpg", "jpeg", "bmp", "tiff", "gif"):
            return _extract_image(file_content)
        elif ext in ("txt", "json", "xml", "md", "log"):
            return file_content.decode("utf-8", errors="ignore")
        else:
            # Tentative de décodage UTF-8 générique
            return file_content.decode("utf-8", errors="ignore")

    except Exception as exc:
        return f"[ERREUR] Impossible d'extraire le texte de {filename}: {exc}"


# ─────────────────────────────────────────────────────────────────────────────
# Extracteurs par format
# ─────────────────────────────────────────────────────────────────────────────

def _extract_pdf(content: bytes, filename: str) -> str:
    if not PdfReader:
        return "[ERREUR] PyPDF2 non installé. Installez avec: pip install pypdf2"

    reader = PdfReader(io.BytesIO(content))
    pages_text = [page.extract_text() or "" for page in reader.pages]
    extracted = "\n\n".join(pages_text).strip()

    # Si aucun texte extrait (PDF scanné), fallback OCR
    if not extracted and Image and pytesseract and convert_from_bytes:
        return _ocr_pdf(content, filename)

    return extracted if extracted else f"[AVERTISSEMENT] Aucun texte extrait de {filename}"


def _ocr_pdf(content: bytes, filename: str) -> str:
    """OCR sur un PDF scanné via pdf2image + Tesseract."""
    try:
        kwargs = {"poppler_path": POPPLER_PATH} if POPPLER_PATH else {}
        images = convert_from_bytes(content, **kwargs)
        parts = []
        for img in images:
            try:
                parts.append(pytesseract.image_to_string(img, lang="fra+eng"))
            except Exception:
                parts.append(pytesseract.image_to_string(img, lang="eng"))
        return "\n\n".join(parts)
    except Exception as exc:
        return f"[ERREUR OCR PDF] {exc}"


def _extract_docx(content: bytes) -> str:
    if not DocxDocument:
        return "[ERREUR] python-docx non installé. Installez avec: pip install python-docx"

    doc = DocxDocument(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _extract_pptx(content: bytes) -> str:
    """Extrait le texte d'une présentation PowerPoint (.pptx / .ppt)."""
    if not PptxPresentation:
        return "[ERREUR] python-pptx non installé. Installez avec: pip install python-pptx"

    try:
        prs = PptxPresentation(io.BytesIO(content))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"=== Diapositive {i} ===")
                parts.append("\n".join(slide_texts))
        return "\n\n".join(parts) if parts else "[AVERTISSEMENT] Aucune texte trouvé dans la présentation."
    except Exception as exc:
        return f"[ERREUR] Impossible d'extraire le texte de la présentation PowerPoint: {exc}"


def _extract_excel(content: bytes) -> str:
    if not pd:
        return "[ERREUR] openpyxl/pandas non installés. Installez avec: pip install openpyxl pandas"

    sheets: dict = pd.read_excel(io.BytesIO(content), sheet_name=None)
    parts = []
    for sheet_name, df in sheets.items():
        parts.append(f"=== Feuille: {sheet_name} ===")
        parts.append(df.to_string(index=False))
    return "\n\n".join(parts)


def _extract_csv(content: bytes) -> str:
    if not pd:
        return content.decode("utf-8", errors="ignore")

    df = pd.read_csv(io.BytesIO(content))
    return df.to_string(index=False)


def _extract_image(content: bytes) -> str:
    if not Image or not pytesseract:
        return "[ERREUR] Pillow/pytesseract non installés. Installez avec: pip install Pillow pytesseract"

    image = Image.open(io.BytesIO(content))
    try:
        text = pytesseract.image_to_string(image, lang="fra+eng")
    except Exception:
        text = pytesseract.image_to_string(image, lang="eng")

    if not text.strip():
        return "[AVERTISSEMENT] Aucun texte détecté dans l'image."
    return text.strip()
